from typing import List
from pathlib import Path
from PyPDF2 import PdfReader
from pptx import Presentation
import markdown
from bs4 import BeautifulSoup
from config.settings import DOCUMENT_PATHS, RAG_CONFIG

class DocumentProcessor:
    def __init__(self):
        self.chunk_size = RAG_CONFIG["chunk_size"]
        self.chunk_overlap = RAG_CONFIG["chunk_overlap"]
    
    def process_all_documents(self) -> List[str]:
        """处理所有文档"""
        texts = []
        
        try:
            # 检查文档目录是否存在
            for path in DOCUMENT_PATHS.values():
                if not path.exists():
                    print(f"警告：文档目录 {path} 不存在，正在创建...")
                    path.mkdir(parents=True, exist_ok=True)
            
            # 处理PDF文件
            pdf_paths = list(DOCUMENT_PATHS["pdfs"].glob("*.pdf"))
            print(f"找到 {len(pdf_paths)} 个PDF文件")
            for pdf_path in pdf_paths:
                print(f"正在处理PDF文件：{pdf_path.name}")
                texts.extend(self._process_pdf(pdf_path))
            
            # 处理PPT文件
            ppt_paths = list(DOCUMENT_PATHS["ppts"].glob("*.pptx"))
            print(f"找到 {len(ppt_paths)} 个PPT文件")
            for ppt_path in ppt_paths:
                print(f"正在处理PPT文件：{ppt_path.name}")
                texts.extend(self._process_ppt(ppt_path))
            
            # 处理Markdown文件
            md_paths = list(DOCUMENT_PATHS["markdown"].glob("*.md"))
            print(f"找到 {len(md_paths)} 个Markdown文件")
            for md_path in md_paths:
                print(f"正在处理Markdown文件：{md_path.name}")
                texts.extend(self._process_markdown(md_path))
            
            if not texts:
                print("警告：没有找到任何文档！请确保在以下目录中放置了文档：")
                for path_type, path in DOCUMENT_PATHS.items():
                    print(f"- {path_type}: {path}")
            
            return texts
            
        except Exception as e:
            print(f"处理文档时出现错误：{str(e)}")
            raise
    
    def _process_pdf(self, file_path: Path) -> List[str]:
        """处理PDF文件"""
        texts = []
        reader = PdfReader(file_path)
        
        for page in reader.pages:
            text = page.extract_text()
            if text:
                texts.extend(self._split_text(text))
        
        return texts
    
    def _process_ppt(self, file_path: Path) -> List[str]:
        """处理PPT文件"""
        try:
            print(f"开始处理PPT文件：{file_path.name}")
            texts = []
            prs = Presentation(file_path)
            
            for i, slide in enumerate(prs.slides, 1):
                print(f"  处理第 {i} 页幻灯片")
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                
                if slide_text:
                    text = "\n".join(slide_text)
                    texts.extend(self._split_text(text))
            
            if not texts:
                print(f"警告：PPT文件 {file_path.name} 中没有提取到任何文本")
            
            return texts
            
        except Exception as e:
            print(f"处理PPT文件 {file_path.name} 时出现错误：{str(e)}")
            raise
    
    def _process_markdown(self, file_path: Path) -> List[str]:
        """处理Markdown文件"""
        texts = []
        with open(file_path, 'r', encoding='utf-8') as f:
            md_content = f.read()
            
        # 将Markdown转换为HTML
        html = markdown.markdown(md_content)
        
        # 使用BeautifulSoup提取纯文本
        soup = BeautifulSoup(html, 'html.parser')
        text = soup.get_text()
        
        if text:
            texts.extend(self._split_text(text))
        
        return texts
    
    def _split_text(self, text: str) -> List[str]:
        """将文本分割成块"""
        if not text or not text.strip():
            return []
            
        chunks = []
        start = 0
        
        while start < len(text):
            end = min(start + self.chunk_size, len(text))
            chunk = text[start:end]
            
            # 如果还没到文本末尾，尝试在句子边界处分割
            if end < len(text):
                # 找到最后一个句号、问号或感叹号
                last_sentence_end = max(
                    chunk.rfind("。"),
                    chunk.rfind("？"),
                    chunk.rfind("！"),
                    chunk.rfind("."),
                    chunk.rfind("?"),
                    chunk.rfind("!")
                )
                
                if last_sentence_end != -1:
                    end = start + last_sentence_end + 1
                    chunk = text[start:end]
            
            chunk = chunk.strip()
            if chunk:  # 只添加非空块
                chunks.append(chunk)
            
            start = end - self.chunk_overlap if end - self.chunk_overlap > start else end
        
        return chunks if chunks else [text.strip()]  # 如果没有分割出块，返回整个文本作为一个块 